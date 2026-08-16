"""
data/ 配下のドキュメント（.pdf / .txt / .md / .docx / .csv / .xlsx / .xls / .pptx / .html / .htm）を
Chroma ベクトルDBに同期するモジュール。

- ライブラリとして: `from ingest import sync_data_dir` を app.py から呼び出し、
  起動時に自動でDBを最新状態に同期する。
- CLIとして: `python ingest.py` を手動実行しても同じ同期処理が走る。

追加・変更されたファイルのみ差分で取り込み、data/ から削除されたファイルは
ベクトルDBからも自動的に削除する（DBとdata/フォルダの内容がズレないようにするため）。
判定には chroma_db/manifest.json にファイル名・更新日時・サイズ・チャンクIDを記録している。

sync_data_dir()はdata/配下を都度全件列挙して差分を検出するため、data/内のファイル数に
比例して処理コストが増える。チャット1往復ごとに会話ログが1ファイルずつ追加される
app.pyの自動ナレッジ化のように「追加対象が1件だけと分かっている」場面向けに、
全件列挙を行わない軽量版の add_single_conversation_file() も用意している。

sync_data_dir()は複数タブ（複数Streamlitセッション）や複数プロセスから同時に呼ばれても
安全なよう、chroma_db/sync.lock を使ったファイルロックで処理全体を排他制御している
（詳細は sync_data_dir() のdocstring参照）。

data/ 直下だけでなく、data/conversations/<thread_id>/ のようなサブフォルダも再帰的に走査する
（app.py のファイルアップロード機能・会話自動保存機能で使用）。会話ログはそのスレッドIDを
チャンクのメタデータ(thread_id)として付与し、rag_chain.build_agent(thread_id) が
「共通ナレッジ＋そのスレッドの会話ログ」だけを検索できるようにしている
（別スレッドの会話ログが回答に混ざらないようにするため）。

PDFの読み込みは2段構成:
  1) PyMuPDF（高速）でまず抽出する。
  2) 抽出できた文字数が極端に少ない場合（図解・スキャンPDFの疑い）だけ、
     Docling（レイアウト認識・OCR内蔵、やや重い）で再解析する。
通常のテキストPDFは1)だけで高速に処理され、図解・スキャンPDFのような
「pypdf/PyMuPDFでは苦手なファイル」だけが2)のコストを払う仕組みにすることで、
処理速度への影響を最小限にしている。Doclingが未インストールの場合は自動的に
1)の結果のみを使う（インストールは任意）。
"""

import argparse
import json
import logging
import os
import re
import sys
from pathlib import Path

from filelock import FileLock, Timeout
from langchain_community.document_loaders import BSHTMLLoader, CSVLoader, Docx2txtLoader, PyMuPDFLoader, TextLoader
from langchain_core.documents import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter

from rag_chain import CHUNK_SIZE, COLLECTION_NAME, GLOBAL_THREAD_ID, PERSIST_DIR, get_vectorstore

try:
    from langchain_docling import DoclingLoader
    from langchain_docling.loader import ExportType

    DOCLING_AVAILABLE = True
except ImportError:
    DOCLING_AVAILABLE = False

logger = logging.getLogger(__name__)


def _log_progress(verbose: bool, message: str, *args) -> None:
    """同期処理の進捗メッセージをloggerに出力する（print()による二重出力を避けるための一本化窓口）。

    verbose=Trueならinfoレベル、Falseならdebugレベルで出力する。読み込み失敗などの警告は
    本関数を使わず、常に見えるべき情報として logger.warning() を直接呼ぶ。
    """
    logger.log(logging.INFO if verbose else logging.DEBUG, message, *args)


DATA_DIR = Path(__file__).parent / "data"
CONVERSATIONS_DIRNAME = "conversations"
MANIFEST_PATH = PERSIST_DIR / "manifest.json"

# 複数タブ・複数プロセス（app.py/api/main.py併用等）からsync_data_dir()が同時に呼ばれた際の
# 排他制御用ロックファイル。manifest読み込み〜DB更新〜manifest書き込みを単一の実行者だけが
# 行うようにし、read-modify-writeの競合（重複登録）を防ぐ。
SYNC_LOCK_PATH = PERSIST_DIR / "sync.lock"
# 通常の同期は数秒〜十数秒で終わるため、それより十分長い待ち時間を設けつつ無限待機は避ける。
SYNC_LOCK_TIMEOUT_SECONDS = 60


class _ExcelLoader:
    """openpyxlでExcel(.xlsx)を読み込む軽量ローダー。

    langchain_communityのUnstructuredExcelLoaderは`unstructured`パッケージ（spacy等を
    含み依存が重い）を必要とするため、openpyxlのみで完結する自前実装にしている。
    シートごとに1つのDocumentを作り、各行をタブ区切りテキストとして連結する。
    """

    def __init__(self, file_path: str):
        self.file_path = file_path

    def load(self) -> list[Document]:
        import openpyxl

        workbook = openpyxl.load_workbook(self.file_path, data_only=True, read_only=True)
        try:
            docs = []
            for sheet in workbook.worksheets:
                lines = [
                    "\t".join("" if cell is None else str(cell) for cell in row)
                    for row in sheet.iter_rows(values_only=True)
                ]
                text = "\n".join(lines).strip()
                if text:
                    docs.append(Document(page_content=text, metadata={"source": self.file_path, "sheet": sheet.title}))
            return docs
        finally:
            # read_only=Trueのワークブックは参照サイクルを持ち、GCが回るまでfdが解放されない
            # ことがあるため、明示的にclose()してsync_data_dir()の全件走査時のfd枯渇を防ぐ。
            workbook.close()


class _LegacyExcelLoader:
    """xlrdで旧形式Excel(.xls)を読み込む軽量ローダー。

    xlrd 2.x以降は.xls専用（.xlsxサポートは廃止済み）で、_ExcelLoaderが使うopenpyxlは
    逆に.xlsを読めないため、拡張子ごとにライブラリを使い分けている。出力形式は
    _ExcelLoaderと揃え、シートごとに1つのDocumentを作りタブ区切りテキストとして連結する。
    """

    def __init__(self, file_path: str):
        self.file_path = file_path

    def load(self) -> list[Document]:
        import xlrd

        workbook = xlrd.open_workbook(self.file_path)
        docs = []
        for sheet in workbook.sheets():
            lines = ["\t".join(str(cell) for cell in sheet.row_values(row_idx)) for row_idx in range(sheet.nrows)]
            text = "\n".join(lines).strip()
            if text:
                docs.append(Document(page_content=text, metadata={"source": self.file_path, "sheet": sheet.name}))
        return docs


class _PowerPointLoader:
    """python-pptxでPowerPoint(.pptx)を読み込む軽量ローダー。

    _ExcelLoaderと同様、依存の重い`unstructured`パッケージを避けるための自前実装。
    スライドごとに1つのDocumentを作り、スライド内のテキストフレームを連結する。
    """

    def __init__(self, file_path: str):
        self.file_path = file_path

    def load(self) -> list[Document]:
        from pptx import Presentation

        presentation = Presentation(self.file_path)
        docs = []
        for index, slide in enumerate(presentation.slides, start=1):
            texts = [
                shape.text_frame.text
                for shape in slide.shapes
                if shape.has_text_frame and shape.text_frame.text.strip()
            ]
            text = "\n".join(texts).strip()
            if text:
                docs.append(Document(page_content=text, metadata={"source": self.file_path, "slide": index}))
        return docs


# 拡張子ごとのローダー対応表（PDFは実際には_load_pdf()で2段構成の判定を行う）
LOADERS = {
    ".pdf": PyMuPDFLoader,
    ".txt": TextLoader,
    ".md": TextLoader,
    ".docx": Docx2txtLoader,
    ".csv": CSVLoader,
    ".xlsx": _ExcelLoader,
    ".xls": _LegacyExcelLoader,
    ".pptx": _PowerPointLoader,
    ".html": BSHTMLLoader,
    ".htm": BSHTMLLoader,
}

# 1ページあたりの抽出文字数がこれ未満の場合、「うまくテキスト抽出できていない
# （図解・スキャンPDFの疑いがある）」とみなしDoclingでの再解析を試みる。
MIN_CHARS_PER_PAGE_FOR_FAST_PATH = 40

# Doclingフォールバック（DOC_CHUNKS）は段落・テーブルセル単位の細切れDocumentを返し、
# そのままsplitterに渡すとチャンクが細かくなりすぎる（splitterはDocumentをまたいで
# マージしない）ため、splitterのchunk_sizeと同じ目安文字数でまとめ直す。
DOCLING_CHUNK_MERGE_TARGET_CHARS = CHUNK_SIZE

# memory.save_conversation() が会話ログのMarkdownに書き込むメタデータ行を検出する正規表現。
FALLBACK_METADATA_PATTERN = re.compile(r"^-\s*一般知識フォールバック:\s*true\s*$", re.MULTILINE)


def safe_upload_dest(filename: str) -> Path | None:
    """アップロードされたファイル名を DATA_DIR 配下の安全な書き込み先パスに変換する。

    ディレクトリ部分（`../` 等）を除いた素のファイル名のみを使い、resolve() 後に
    DATA_DIR 配下から外れていないか（パストラバーサルの疑いがないか）を最終チェックする。
    外れる場合は None を返す。同名ファイルの重複チェックは呼び出し元の責務。
    """
    dest = (DATA_DIR / Path(filename).name).resolve()
    if dest.parent != DATA_DIR.resolve():
        return None
    return dest


def safe_relative_dest(relative_path: str) -> Path | None:
    """既知の相対パス（サブフォルダを含みうる）を DATA_DIR 配下の安全な実パスに変換する。

    safe_upload_dest() と異なりディレクトリ部分を切り捨てない（サブフォルダ構造を
    保ったまま対象ファイルを一意に特定する）。DATA_DIR 配下から外れる場合
    （`../` によるパストラバーサルの疑いがある場合）は None を返す。
    """
    data_dir_resolved = DATA_DIR.resolve()
    dest = (DATA_DIR / relative_path).resolve()
    if dest != data_dir_resolved and data_dir_resolved not in dest.parents:
        return None
    return dest


def resolve_upload_dest(filename: str, taken_paths: set[Path] | None = None) -> Path | None:
    """アップロードされたファイル名から、上書きを避けた実際の書き込み先パスを求める。

    同名ファイルが既に存在する場合（taken_pathsも含め同一バッチ内の重複も対象）、
    無警告での上書きを避けるため "name (2).ext" のように連番サフィックスを付けた
    空いているパスを返す。呼び出し元は戻り値のファイル名が元と異なる場合に警告表示する
    想定。パストラバーサルの疑いがある場合は safe_upload_dest() と同様に None を返す。
    """
    dest = safe_upload_dest(filename)
    if dest is None:
        return None

    taken = taken_paths if taken_paths is not None else set()
    if dest not in taken and not dest.exists():
        return dest

    stem, suffix = dest.stem, dest.suffix
    counter = 2
    candidate = dest.with_name(f"{stem} ({counter}){suffix}")
    while candidate in taken or candidate.exists():
        counter += 1
        candidate = dest.with_name(f"{stem} ({counter}){suffix}")
    return candidate


def _load_manifest() -> dict:
    if not MANIFEST_PATH.exists():
        return {}
    try:
        return json.loads(MANIFEST_PATH.read_text(encoding="utf-8"))
    except json.JSONDecodeError:
        logger.warning(
            "%s の読み込みに失敗しました（壊れたJSONの可能性があります）。"
            "空のマニフェストとして扱い、全ファイルを再取り込みします。",
            MANIFEST_PATH,
        )
        return {}
    except FileNotFoundError:
        # exists()での確認直後にファイルが削除された場合（TOCTOU）に備える。
        logger.warning(
            "%s の読み込み中にファイルが見つかりませんでした。"
            "空のマニフェストとして扱い、全ファイルを再取り込みします。",
            MANIFEST_PATH,
        )
        return {}


def _save_manifest(manifest: dict) -> None:
    PERSIST_DIR.mkdir(parents=True, exist_ok=True)
    tmp_path = MANIFEST_PATH.with_suffix(".json.tmp")
    tmp_path.write_text(json.dumps(manifest, ensure_ascii=False, indent=2), encoding="utf-8")
    os.replace(tmp_path, MANIFEST_PATH)


def _fingerprint(path: Path) -> dict:
    stat = path.stat()
    return {"mtime": stat.st_mtime, "size": stat.st_size}


def _load_pdf(path: Path, verbose: bool = True) -> list:
    """PDFを読み込む（PyMuPDFで高速抽出 → 必要な場合のみDoclingでフォールバック）。"""
    try:
        fast_docs = PyMuPDFLoader(str(path)).load()
    except Exception as e:
        # PyMuPDF自体が例外を送出した場合（暗号化PDF・破損PDF・特殊なPDF構造など）。
        # Doclingが利用可能なら、レイアウト認識・OCRで読める可能性があるためフォールバックする。
        # Doclingが未インストール、またはDoclingも失敗した場合は元の例外をそのまま送出し、
        # 呼び出し元（sync_data_dir）で "failed" として記録・次回リトライさせる。
        if not DOCLING_AVAILABLE:
            raise
        _log_progress(verbose, "    → PyMuPDFでの読み込みに失敗したため、Doclingでの再解析を試みます（%s）...", e)
        docling_docs = _load_pdf_with_docling(path, verbose=verbose)
        if not docling_docs:
            raise
        docling_chars = sum(len(d.page_content.strip()) for d in docling_docs)
        _log_progress(verbose, "    → Doclingで%d文字を抽出しました。", docling_chars)
        return docling_docs

    total_chars = sum(len(d.page_content.strip()) for d in fast_docs)
    avg_chars_per_page = total_chars / max(len(fast_docs), 1)

    if avg_chars_per_page >= MIN_CHARS_PER_PAGE_FOR_FAST_PATH or not DOCLING_AVAILABLE:
        return fast_docs

    _log_progress(
        verbose,
        "    → テキスト抽出量が少ない（平均%.0f文字/ページ）ため、"
        "Doclingで図解・OCR解析を試みます（時間がかかる場合があります）...",
        avg_chars_per_page,
    )
    docling_docs = _load_pdf_with_docling(path, verbose=verbose)
    if docling_docs:
        docling_chars = sum(len(d.page_content.strip()) for d in docling_docs)
        if docling_chars > total_chars:
            _log_progress(
                verbose, "    → Doclingで%d文字を抽出しました（PyMuPDF: %d文字）。", docling_chars, total_chars
            )
            return docling_docs

    return fast_docs


def _load_pdf_with_docling(path: Path, verbose: bool = True) -> list:
    """Doclingでの再解析を試みる。失敗した場合は空リストを返す（例外は送出しない）。

    export_type=DOC_CHUNKS を使うのは、PyMuPDFLoaderと同様にページ番号相当のメタデータ
    （page）を付与できるため（export_type=MARKDOWNではページ境界の情報が失われる）。
    DOC_CHUNKSが返す個々のDocumentは段落・テーブルセル単位と細かいため、後段の
    splitterに渡す前に_merge_docling_chunks()である程度まとめておく。
    """
    try:
        docling_docs = DoclingLoader(file_path=str(path), export_type=ExportType.DOC_CHUNKS).load()
    except Exception as e:
        _log_progress(verbose, "    → Docling解析に失敗しました（%s）", e)
        return []
    for doc in docling_docs:
        doc.metadata.setdefault("source", str(path))
    return _merge_docling_chunks(docling_docs)


def _merge_docling_chunks(docling_docs: list, target_chars: int = DOCLING_CHUNK_MERGE_TARGET_CHARS) -> list[Document]:
    """DOC_CHUNKSが返す細切れのDocumentを、target_chars程度になるまで隣接結合する。

    構造的な区切りは考慮せず文字数のみで区切る単純な方式（最終的な文単位の分割は
    後段のRecursiveCharacterTextSplitterが担うため、ここでは粒度を揃えるだけでよい）。
    各グループのpageメタデータは、含まれる各Documentの元ページ番号の最小値を採用し、
    有効なページ番号が1つも無い場合は付与しない。
    """
    merged_docs: list[Document] = []
    source = docling_docs[0].metadata.get("source") if docling_docs else None
    buffer_texts: list[str] = []
    buffer_pages: list[int] = []
    buffer_len = 0

    def flush() -> None:
        if not buffer_texts:
            return
        metadata = {}
        if source is not None:
            metadata["source"] = source
        if buffer_pages:
            metadata["page"] = min(buffer_pages)
        merged_docs.append(Document(page_content="\n\n".join(buffer_texts), metadata=metadata))

    for doc in docling_docs:
        text = doc.page_content
        page = _extract_docling_page(doc.metadata)
        if buffer_texts and buffer_len + len(text) > target_chars:
            flush()
            buffer_texts, buffer_pages, buffer_len = [], [], 0
        buffer_texts.append(text)
        if page is not None:
            buffer_pages.append(page)
        buffer_len += len(text)

    flush()
    return merged_docs


def _extract_docling_page(metadata: dict) -> int | None:
    """DoclingLoader(export_type=DOC_CHUNKS)が付与するdl_metaから代表ページ番号を取り出す。

    1チャンクが複数ページにまたがりうるため、表示用に最小のpage_no（最初のページ）を採用する。
    Doclingのpage_noは1始まりのため、PyMuPDFLoaderが付与するpage（0始まり）と揃うよう
    1引いて返す。想定外の形式・値の場合はNoneを返す。
    """
    dl_meta = metadata.get("dl_meta")
    if not isinstance(dl_meta, dict):
        return None
    page_numbers = [
        prov["page_no"]
        for item in dl_meta.get("doc_items", [])
        if isinstance(item, dict)
        for prov in item.get("prov", [])
        if isinstance(prov, dict) and isinstance(prov.get("page_no"), int) and prov["page_no"] >= 1
    ]
    if not page_numbers:
        return None
    return min(page_numbers) - 1


def _thread_id_for(rel_path: str) -> str:
    """相対パスから会話スレッドIDを判定する。

    data/conversations/<thread_id>/xxx.md → そのthread_id（このスレッドでのみ検索対象）
    data/conversations/xxx.md（サブフォルダなし、スレッド機能追加前の古い保存形式）
        → GLOBAL_THREAD_ID として扱う（後方互換。以前の会話ログも消えず、共通ナレッジとして生き続ける）
    それ以外（data/直下のファイルやアップロードファイルなど） → GLOBAL_THREAD_ID（全スレッド共通）
    """
    parts = Path(rel_path).parts
    if len(parts) >= 3 and parts[0] == CONVERSATIONS_DIRNAME:
        return parts[1]
    return GLOBAL_THREAD_ID


def _is_fallback_conversation(docs: list) -> bool:
    """分割前の生ドキュメントの本文から、一般知識フォールバック回答の会話ログかどうかを判定する。

    memory.save_conversation() が書き込む「- 一般知識フォールバック: true」という
    メタデータ行を正規表現で検出する。
    """
    return any(FALLBACK_METADATA_PATTERN.search(doc.page_content) for doc in docs)


def data_dir_signature() -> tuple[int, float]:
    """data/ の変更有無を、内容を読まずにstat()だけで軽量に判定するためのシグネチャを返す。

    sync_data_dir()が対象とするのと同じファイル集合に対し (ファイル数, 最新mtime) の
    タプルを返す（あくまで近似的な変更検知で、読み込み・埋め込み等の重い処理は行わない）。
    app.py側は再実行のたびにこれを前回値と比較し、変更があった場合だけ本格的な
    sync_data_dir() を呼び出す。対象ファイルが1つも無い場合は (0, 0.0) を返す。
    """
    if not DATA_DIR.exists():
        return (0, 0.0)
    target_files = [f for f in DATA_DIR.rglob("*") if f.is_file() and f.suffix.lower() in LOADERS]
    if not target_files:
        return (0, 0.0)
    latest_mtime = max(f.stat().st_mtime for f in target_files)
    return (len(target_files), latest_mtime)


def sync_data_dir(verbose: bool = True) -> dict:
    """data/ の内容とベクトルDBを同期する。

    戻り値: {"added": [...], "updated": [...], "removed": [...], "failed": [...]}
    （変更がなければ全て空リスト）。読み込み・分割に失敗したファイルは"failed"に積まれ
    manifestには記録されない（次回同期時にリトライされる。他ファイルの同期は継続する）。

    manifestはプロセス中断時の重複登録を防ぐためファイル1件ごとに保存する。また、
    既存ファイル更新時の旧チャンクdelete()が失敗した場合は旧chunk_idsを
    pending_delete_chunk_idsとしてmanifestに持ち越し、以降unchanged判定になっても
    削除だけは再試行する（さもないと新旧チャンクが重複したまま残り続けるため）。

    複数タブ・複数プロセスから同時に呼ばれても安全なよう、manifest読み込み〜DB更新〜
    manifest書き込みをファイルロック（SYNC_LOCK_PATH）で排他制御する。
    SYNC_LOCK_TIMEOUT_SECONDS以内にロックを獲得できなかった場合は
    `filelock.Timeout` をそのまま送出する。
    """
    DATA_DIR.mkdir(parents=True, exist_ok=True)
    PERSIST_DIR.mkdir(parents=True, exist_ok=True)

    try:
        with FileLock(str(SYNC_LOCK_PATH), timeout=SYNC_LOCK_TIMEOUT_SECONDS):
            return _sync_data_dir_locked(verbose=verbose)
    except Timeout:
        logger.warning(
            "%s のロック取得がタイムアウトしました（他のセッションが同期中の可能性があります）。",
            SYNC_LOCK_PATH,
        )
        raise


def add_single_conversation_file(path: Path) -> str:
    """会話ログ1件だけを、data/ 全件を走査せずにその場でベクトルDBへ反映する軽量な経路。

    sync_data_dir()はDATA_DIR.rglob("*")による全件列挙を伴うため、チャット1往復ごとに
    毎回呼ぶとdata/配下のファイル数に比例して処理コストが増え続けてしまう。本関数は
    対象ファイル1件分の読み込み〜manifest更新だけを行い、全件列挙は行わない。

    sync_data_dir()と同じSYNC_LOCK_PATHのファイルロックで排他制御する。

    戻り値: "added" / "updated" / "unchanged" / "failed" のいずれか
    （_ingest_file()の戻り値をそのまま返す）。ロック取得がタイムアウトした場合は
    `filelock.Timeout` をそのまま送出する。
    """
    DATA_DIR.mkdir(parents=True, exist_ok=True)
    PERSIST_DIR.mkdir(parents=True, exist_ok=True)

    try:
        with FileLock(str(SYNC_LOCK_PATH), timeout=SYNC_LOCK_TIMEOUT_SECONDS):
            return _add_single_conversation_file_locked(path)
    except Timeout:
        logger.warning(
            "%s のロック取得がタイムアウトしました（他のセッションが同期中の可能性があります）。",
            SYNC_LOCK_PATH,
        )
        raise


def _add_single_conversation_file_locked(path: Path) -> str:
    """add_single_conversation_file()の本体（呼び出し元がファイルロックを取得済みであることが前提）。"""
    vector_store = get_vectorstore()
    manifest = _load_manifest()
    splitter = RecursiveCharacterTextSplitter(chunk_size=CHUNK_SIZE, chunk_overlap=200)

    name = str(path.relative_to(DATA_DIR))
    status = _ingest_file(name, path, vector_store, manifest, splitter, verbose=False)
    # unchanged・failedの場合は_ingest_file()内で保存されないため、ここで冪等に保存しておく。
    _save_manifest(manifest)
    return status


def _ingest_file(name: str, path: Path, vector_store, manifest: dict, splitter, verbose: bool) -> str:
    """1ファイル分の追加・更新判定と、必要な場合のベクトルDBへの反映を行う。

    sync_data_dir()の全件差分検出、add_single_conversation_file()の単一ファイル追加の
    両方から呼ばれる共通処理。manifestの読み込み自体は呼び出し元の責務とし、
    ここでは受け取ったmanifest辞書をその場で書き換える。

    戻り値: "added" / "updated" / "unchanged" / "failed" のいずれか。
    """
    fingerprint = _fingerprint(path)
    entry = manifest.get(name)
    unchanged = entry and entry.get("mtime") == fingerprint["mtime"] and entry.get("size") == fingerprint["size"]
    if unchanged:
        # 内容に変化がなくても、前回の旧チャンク削除が失敗し持ち越しになっている場合はここで再試行する。
        pending_delete_chunk_ids = entry.get("pending_delete_chunk_ids")
        if pending_delete_chunk_ids:
            try:
                vector_store.delete(ids=pending_delete_chunk_ids)
            except Exception as e:
                logger.warning(
                    "%s の保留中だった旧チャンク削除の再試行に失敗しました（次回同期時に再試行します）: %s",
                    name,
                    e,
                )
            else:
                del entry["pending_delete_chunk_ids"]
                _save_manifest(manifest)
                _log_progress(verbose, "%s: 保留中だった旧チャンクの削除が完了しました。", name)
        return "unchanged"

    try:
        if path.suffix.lower() == ".pdf":
            docs = _load_pdf(path, verbose=verbose)
        else:
            loader = LOADERS[path.suffix.lower()](str(path))
            docs = loader.load()
        chunks = splitter.split_documents(docs)
    except Exception as e:
        # 1ファイルの読み込み失敗で他の正常なファイルの同期まで止めないよう、
        # ログに残してスキップする（manifestには記録しないので次回リトライされる）。
        logger.warning("%s の読み込みに失敗したためスキップします: %s", name, e)
        return "failed"

    thread_id = _thread_id_for(name)
    # 根拠のない一般知識フォールバック回答がそのまま再学習され「裏付けのある回答」として
    # 再ヒットしないよう、全チャンクにis_fallbackを付与しretrieve_context側で除外させる。
    is_fallback = _is_fallback_conversation(docs)
    for chunk in chunks:
        chunk.metadata["thread_id"] = thread_id
        chunk.metadata["is_fallback"] = is_fallback

    # 新チャンクの追加を先に行い、成功してから旧チャンクを削除する。delete→addの順だと
    # 追加失敗時に「旧チャンクが消えたのに新チャンクも無い」データ消失状態になりうるため、
    # add→deleteの逆順にして追加失敗時は旧チャンクが残る方（安全側）に倒す。
    try:
        chunk_ids = vector_store.add_documents(documents=chunks) if chunks else []
    except Exception as e:
        logger.warning("%s のベクトルストアへの追加に失敗したためスキップします: %s", name, e)
        return "failed"

    # 旧チャンクのdelete()が失敗した場合、そのまま握りつぶすとmanifestが新内容で
    # 更新されて次回unchanged判定になり、削除が二度と再試行されず重複が残り続ける。
    # 失敗時は旧chunk_idsをpending_delete_chunk_idsとして持ち越し、unchanged判定時に
    # 再試行する（前回分の持ち越しが残っていれば今回分と合算する）。
    pending_delete_chunk_ids = set(entry.get("pending_delete_chunk_ids") or []) if entry else set()
    if entry and entry.get("chunk_ids"):
        try:
            vector_store.delete(ids=entry["chunk_ids"])
        except Exception as e:
            logger.warning("%s の旧チャンク削除に失敗しました（次回同期時に再試行します）: %s", name, e)
            pending_delete_chunk_ids |= set(entry["chunk_ids"])

    manifest[name] = {**fingerprint, "chunk_ids": chunk_ids}
    if pending_delete_chunk_ids:
        manifest[name]["pending_delete_chunk_ids"] = sorted(pending_delete_chunk_ids)
    status = "updated" if entry else "added"
    action = "更新" if entry else "追加"
    _log_progress(verbose, "%s: %s（%dチャンク）", action, name, len(chunks))

    # 全件処理後にまとめて保存すると、DB追加後・保存前にプロセスが中断した場合に
    # DBには反映済みなのにmanifest未記録の不整合（重複登録の原因）が起きるため都度保存する。
    _save_manifest(manifest)
    return status


def _sync_data_dir_locked(verbose: bool) -> dict:
    """sync_data_dir()の本体（呼び出し元がファイルロックを取得済みであることが前提）。"""
    vector_store = get_vectorstore()
    manifest = _load_manifest()
    splitter = RecursiveCharacterTextSplitter(chunk_size=CHUNK_SIZE, chunk_overlap=200)

    # data/ 直下だけでなく、data/conversations/ などのサブフォルダも再帰的に走査する
    current_files = {
        str(f.relative_to(DATA_DIR)): f for f in DATA_DIR.rglob("*") if f.is_file() and f.suffix.lower() in LOADERS
    }

    result = {"added": [], "updated": [], "removed": [], "failed": []}

    for name, path in current_files.items():
        status = _ingest_file(name, path, vector_store, manifest, splitter, verbose=verbose)
        if status in ("added", "updated", "failed"):
            result[status].append(name)

    for name in list(manifest.keys()):
        if name not in current_files:
            entry = manifest[name]
            # manifestエントリごと消すと再試行の機会を失うため、保留中の旧チャンクも
            # 現行のchunk_idsと合わせて削除対象にする。
            ids_to_delete = set(entry.get("chunk_ids") or []) | set(entry.get("pending_delete_chunk_ids") or [])
            if ids_to_delete:
                try:
                    vector_store.delete(ids=sorted(ids_to_delete))
                except Exception as e:
                    # 削除失敗時はpending_delete_chunk_idsだけ残し、次回同期時にこのループで再試行する。
                    logger.warning(
                        "%s の削除処理中に旧チャンク削除に失敗しました（次回同期時に再試行します）: %s",
                        name,
                        e,
                    )
                    manifest[name] = {"pending_delete_chunk_ids": sorted(ids_to_delete)}
                    _save_manifest(manifest)
                    continue
            del manifest[name]
            result["removed"].append(name)
            _log_progress(verbose, "削除: %s", name)
            # 追加・更新時と同様、削除もDB反映とmanifest保存を1件ごとに一致させる。
            _save_manifest(manifest)

    # 全ファイルが失敗し1件も追加・更新・削除が無かった場合でも、manifest.jsonが
    # 存在しない状態のまま終わらないよう最後に保存しておく（内容は変わらないため冪等）。
    _save_manifest(manifest)

    if not any(result.values()):
        _log_progress(verbose, "変更はありませんでした（すでに同期済みです）。")

    return result


def list_indexed_files() -> list[dict]:
    """サイドバーの一覧表示用に、インデックス済みファイルの情報をmanifestから取得する。

    会話ログ（data/conversations/配下）は除外し、data/直下のドキュメントのみを返す。
    戻り値はファイル名昇順のリストで、各要素は {"name": ファイル名, "chunk_count": チャンク数}。
    delete()失敗時にpending_delete_chunk_idsだけを持ち越した「ゴーストエントリ」
    （"chunk_ids"キーを持たないエントリ）は内部状態のため一覧には含めない。
    """
    manifest = _load_manifest()
    return [
        {"name": name, "chunk_count": len(entry["chunk_ids"])}
        for name, entry in sorted(manifest.items())
        if Path(name).parts[0] != CONVERSATIONS_DIRNAME and "chunk_ids" in entry
    ]


def delete_indexed_file(name: str) -> bool:
    """インデックス済みファイルをdata/から削除する。

    ここではmanifest・ベクトルDBの更新は行わない。呼び出し元がこの後sync_data_dir()を
    呼ぶことで、ファイル消失が検知されDB・manifestから自動的に除外される。
    相対パスはサブフォルダ（例: "manuals/spec.pdf"）を含みうるため、safe_upload_dest()
    ではなくsafe_relative_dest()でパスを解決する。対象が存在しない場合はFalseを返す。
    """
    path = safe_relative_dest(name)
    if path is None or not path.exists():
        return False
    path.unlink()
    return True


def _dir_size_mb(path: Path) -> float:
    if not path.exists():
        return 0.0
    total = sum(f.stat().st_size for f in path.rglob("*") if f.is_file())
    return total / (1024 * 1024)


def print_status() -> None:
    """ベクトルDBの現在の状態を表示する（`python ingest.py --status`）。"""
    manifest = _load_manifest()

    print(f"DB保存先     : {PERSIST_DIR}")
    print(f"コレクション名: {COLLECTION_NAME}")

    if not PERSIST_DIR.exists() or not manifest:
        print("インデックス済みのファイルはまだありません。")
        print("`streamlit run app.py` を起動するか `python ingest.py` を実行してください。")
        return

    try:
        vector_store = get_vectorstore()
        chunk_count = vector_store._collection.count()
    except Exception:
        chunk_count = sum(len(v.get("chunk_ids", [])) for v in manifest.values())

    print(f"インデックス済みファイル数: {len(manifest)}件")
    print(f"チャンク数（ベクトル数）  : {chunk_count}件")
    print(f"DBフォルダのサイズ       : {_dir_size_mb(PERSIST_DIR):.1f} MB")
    print("\nファイル別チャンク数:")
    for name, entry in sorted(manifest.items()):
        print(f"  - {name}: {len(entry.get('chunk_ids', []))}チャンク")


def main():
    parser = argparse.ArgumentParser(description="data/ とベクトルDBの同期・状態確認ツール")
    parser.add_argument(
        "--status",
        action="store_true",
        help="同期は行わず、現在のDBの状態（件数・容量など）だけを表示する",
    )
    args = parser.parse_args()

    if args.status:
        print_status()
        return

    # sync_data_dir()の進捗ログ（logger.info）をCLI実行時にコンソールへ表示するための設定。
    # stream=sys.stdoutにするのは、同じmain()内のprint()出力と表示順序を揃えるため
    # （logging標準のデフォルトはstderrで、print()と混在すると順序が入れ替わりうる）。
    logging.basicConfig(level=logging.INFO, format="%(message)s", stream=sys.stdout)

    print("data/ をベクトルDBに同期しています...")
    result = sync_data_dir()
    print(
        f"完了: 追加{len(result['added'])}件 / "
        f"更新{len(result['updated'])}件 / 削除{len(result['removed'])}件 / "
        f"失敗{len(result['failed'])}件"
    )
    print("`streamlit run app.py` でチャットを開始できます。")
    print("（DBの状態を確認したい場合は `python ingest.py --status`）")


if __name__ == "__main__":
    main()
